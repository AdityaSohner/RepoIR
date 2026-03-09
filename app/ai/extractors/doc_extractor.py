import os
import io
import pdfplumber
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
from pdf2image import convert_from_path, convert_from_bytes
import pytesseract
from .base import BaseExtractor

class DocumentExtractor(BaseExtractor):
    def extract(self, source, extension=None) -> dict:
        if extension is None:
            if isinstance(source, str):
                extension = os.path.splitext(source)[1].lower()
            else:
                raise ValueError("Extension must be provided for stream sources")

        if extension == ".pdf":
            text = self._extract_pdf(source)
        elif extension == ".docx":
            text = self._extract_docx(source)
        elif extension == ".pptx":
            text = self._extract_pptx(source)
        elif extension in [".xlsx", ".xls"]:
            text = self._extract_excel(source)
        elif extension == ".txt":
            text = self._extract_txt(source)
        else:
            raise ValueError(f"Unsupported document type: {extension}")

        return {
            "raw_text": text.strip(),
            "metadata": {
                "file_type": extension
            }
        }

    def _extract_pdf(self, source) -> str:
        text_parts = []
        with pdfplumber.open(source) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)

        if isinstance(source, (io.BytesIO, bytes)):
            if hasattr(source, 'seek'): source.seek(0)
            data = source.read() if hasattr(source, 'read') else source
            images = convert_from_bytes(data)
        else:
            images = convert_from_path(source)

        for img in images:
            ocr_text = pytesseract.image_to_string(img)
            if ocr_text.strip():
                text_parts.append(ocr_text)
        return "\n".join(text_parts)

    def _extract_docx(self, source) -> str:
        doc = Document(source)
        text_parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                text_parts.append(p.text)
        for rel in doc.part._rels.values():
            if "image" in rel.target_ref:
                image_bytes = rel.target_part.blob
                image = Image.open(io.BytesIO(image_bytes))
                ocr_text = pytesseract.image_to_string(image)
                if ocr_text.strip():
                    text_parts.append(ocr_text)
        return "\n".join(text_parts)

    def _extract_pptx(self, source) -> str:
        prs = Presentation(source)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        return "\n".join(text_parts)

    def _extract_excel(self, source) -> str:
        wb = load_workbook(source, data_only=True)
        summaries = []
        for sheet in wb.worksheets:
            headers = []
            for row in sheet.iter_rows(min_row=1, max_row=1, values_only=True):
                headers = [str(c) for c in row if c]
            if headers:
                summaries.append(f"sheet {sheet.title} with columns " + ", ".join(headers))
        return ". ".join(summaries)

    def _extract_txt(self, source) -> str:
        if hasattr(source, 'read'):
            if hasattr(source, 'seek'): source.seek(0)
            return source.read().decode('utf-8', errors='ignore')
        with open(source, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
